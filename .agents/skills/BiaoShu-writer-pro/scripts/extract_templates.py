#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
标书模板知识提取脚本
从标书范本文件中提取内容，转换为Markdown格式
版本：v1.0 (2026-05-21)
"""

import io
import os
import sys
from pathlib import Path

# 设置标准输出编码（解决Windows控制台编码问题）
import locale
try:
    # 尝试设置控制台编码
    if hasattr(sys.stdout, 'reconfigure'):
        sys.stdout.reconfigure(encoding='utf-8')
    if hasattr(sys.stderr, 'reconfigure'):
        sys.stderr.reconfigure(encoding='utf-8')
except Exception:
    pass


def extract_docx_content(file_path: str) -> str:
    """从docx文件提取内容"""
    try:
        from docx import Document
        doc = Document(file_path)
        content = []

        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                # 检测标题样式
                style_name = para.style.name.lower()
                if 'heading 1' in style_name:
                    content.append(f'# {text}')
                elif 'heading 2' in style_name:
                    content.append(f'## {text}')
                elif 'heading 3' in style_name:
                    content.append(f'### {text}')
                else:
                    content.append(text)

        return '\n\n'.join(content)
    except Exception as e:
        return f"[错误] 无法读取 {file_path}: {str(e)}"


def extract_doc_content(file_path: str) -> str:
    """从doc文件提取内容（需要pywin32）"""
    try:
        import win32com.client
        import pythoncom

        pythoncom.CoInitialize()
        word = win32com.client.Dispatch("Word.Application")
        word.Visible = False

        try:
            doc = word.Documents.Open(os.path.abspath(file_path))
            content = doc.Content.Text
            doc.Close(False)
            return content
        finally:
            word.Quit()
            pythoncom.CoUninitialize()
    except ImportError:
        return "[跳过] pywin32未安装，无法读取.doc文件"
    except Exception as e:
        return f"[错误] 无法读取 {file_path}: {str(e)}"


def extract_pptx_content(file_path: str) -> str:
    """从pptx文件提取内容"""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        content = []

        for i, slide in enumerate(prs.slides, 1):
            content.append(f'## 幻灯片 {i}')
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    content.append(shape.text.strip())

        return '\n\n'.join(content)
    except ImportError:
        return "[跳过] python-pptx未安装，无法读取.pptx文件"
    except Exception as e:
        return f"[错误] 无法读取 {file_path}: {str(e)}"


def process_templates_folder(source_dir: str, output_dir: str):
    """处理模板文件夹"""
    os.makedirs(output_dir, exist_ok=True)

    # 知识文档列表（按重要性排序）
    knowledge_docs = [
        '投标书制作-及范本讲解.doc',
        '如何编好投标文件.docx',
        '商务标、技术标、经济标的定义与区别.doc',
        '投标书版本、格式.doc',
        '投标书格式(范本).doc',
        '投标书格式范本.docx',
        'FS常用投标报价8大技巧.doc',
    ]

    # 范本文档列表
    template_docs = [
        '投标书模版  (完整版).docx',
        '投标书模版--(完整版).docx',
        '投标文件范本.doc',
        '投标文件范本(定稿).doc',
        '投标文件范本(正规).doc',
        '投标文件(范本).docx',
        '投标函(格式).docx',
        '投标函范本.docx',
        'XX公司投标书范本.doc',
        'XX公司投标书范本.docx',
    ]

    # 处理知识文档
    print("=" * 60)
    print("处理知识文档...")
    print("=" * 60)

    for doc_name in knowledge_docs:
        doc_path = os.path.join(source_dir, doc_name)
        if os.path.exists(doc_path):
            print(f"\n处理: {doc_name}")

            if doc_name.endswith('.docx'):
                content = extract_docx_content(doc_path)
            elif doc_name.endswith('.doc'):
                content = extract_doc_content(doc_path)
            else:
                continue

            # 保存为md文件
            md_name = doc_name.rsplit('.', 1)[0] + '.md'
            md_path = os.path.join(output_dir, f"knowledge_{md_name}")

            with open(md_path, 'w', encoding='utf-8') as f:
                f.write(f"# {doc_name}\n\n")
                f.write(content)

            print(f"  已保存: {md_path}")

    # 处理范本文档
    print("\n" + "=" * 60)
    print("处理范本文档...")
    print("=" * 60)

    for doc_name in template_docs:
        doc_path = os.path.join(source_dir, doc_name)
        if os.path.exists(doc_path):
            print(f"\n处理: {doc_name}")

            if doc_name.endswith('.docx'):
                content = extract_docx_content(doc_path)
            elif doc_name.endswith('.doc'):
                content = extract_doc_content(doc_path)
            else:
                continue

            # 保存为md文件
            md_name = doc_name.rsplit('.', 1)[0] + '.md'
            md_path = os.path.join(output_dir, f"template_{md_name}")

            with open(md_path, 'w', encoding='utf-8') as f:
                f.write(f"# {doc_name}\n\n")
                f.write(content)

            print(f"  已保存: {md_path}")

    print("\n" + "=" * 60)
    print("处理完成！")
    print("=" * 60)


def main():
    import argparse

    parser = argparse.ArgumentParser(description='标书模板知识提取工具')
    parser.add_argument('--source', '-s',
                       default='./source_templates',
                       help='源文件夹路径')
    parser.add_argument('--output', '-o',
                       default='./templates',
                       help='输出文件夹路径')

    args = parser.parse_args()

    if not os.path.exists(args.source):
        print(f"[ERROR] 源文件夹不存在: {args.source}")
        sys.exit(1)

    process_templates_folder(args.source, args.output)


if __name__ == '__main__':
    main()
